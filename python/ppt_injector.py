import json
import os
from os import path

import pandas as pd
from pptx import Presentation

THUMBNAIL = "static"

def replace_tokens_and_links(shape, token_map, link_map):
    """
    Replace token placeholders with values and attach hyperlinks in a PowerPoint shape.
    
    Constraint: Tokens must exist fully within a single run. PowerPoint may split
    text across runs if template was edited with font changes or copy-pasted from Word.
    If a token is split across runs, it will not be replaced.
    
    Args:
        shape: PowerPoint shape object
        token_map: Dictionary mapping tokens to replacement display text
        link_map: Dictionary mapping tokens to URLs
    """
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for token, text_value in token_map.items():
                if token in run.text:
                    run.text = run.text.replace(token, text_value)
                    
                    if token in link_map and link_map[token]:
                        run.hyperlink.address = link_map[token]
                    break


def inject_kitchen_snapshot(prs, player_id, graphics_dir):
    """
    Inject kitchen snapshot image into the PowerPoint presentation.
    
    Looks for the slide titled "NethriQ Insight 3 – Kitchen Control" and
    inserts the kitchen snapshot image into the placeholder named "KITCHEN_SNAPSHOT".
    
    Args:
        prs: PowerPoint presentation object
        player_id: Player ID (used to construct the image filename)
        graphics_dir: Directory containing the generated PNG images
    """
    image_filename = f"kitchen_player_{int(player_id)}.png"
    image_path = path.join(graphics_dir, image_filename)
    
    if not path.exists(image_path):
        print(f"Warning: Image not found at {image_path}")
        return
    
    # Find the slide with title "NethriQ Insight 3 – Kitchen Control"
    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_title = shape.text.strip()
                if "NethriQ Insight 3 – Kitchen Control" in slide_title or "NethriQ Insight 3 - Kitchen Control" in slide_title:
                    target_slide = slide
                    break
        if target_slide:
            break
    
    if not target_slide:
        print("Warning: Slide titled 'NethriQ Insight 3 – Kitchen Control' not found")
        return
    
    # Find and replace the KITCHEN_SNAPSHOT placeholder
    for shape in target_slide.shapes:
        if shape.name == "KITCHEN_SNAPSHOT":
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # Remove placeholder
            shape._element.getparent().remove(shape._element)

            # Add image
            target_slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=width,
                height=height
            )
            return
    
    print("Warning: Placeholder 'KITCHEN_SNAPSHOT' not found in the target slide")

def inject_thumbnail(prs, player_id, media_dir):
    """
    Inject hero thumbnail into the slide titled "NethriQ Benchmarks" and
    replace the image placeholder named "THUMBNAIL".
    """
    player_key = f"player_{int(player_id)}"
    candidate_path = path.join(media_dir, "players", player_key, "hero", "hero_thumbnail.jpg")
    image_path = candidate_path if path.exists(candidate_path) else None

    if not image_path:
        print(f"Warning: Thumbnail not found for {player_key}")
        return

    target_slide = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_title = shape.text.strip()
                if "NethriQ Benchmarks" in slide_title:
                    target_slide = slide
                    break
        if target_slide:
            break

    if not target_slide:
        print("Warning: Slide titled 'NethriQ Benchmarks' not found")
        return

    for shape in target_slide.shapes:
        if shape.name == "THUMBNAIL":
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            shape._element.getparent().remove(shape._element)

            target_slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=width,
                height=height
            )
            return

    print("Warning: Placeholder 'THUMBNAIL' not found in the target slide")

# ============================================================================
# Set up file paths and load data
# ============================================================================
csv_path = path.join(path.dirname(path.abspath(__file__)), '..', 'data', 'player_data', 'player_averages.csv')
ppt_template_path = path.join(path.dirname(path.abspath(__file__)),'..', 'node', 'mixed_doubles', 'NethriQ_Gautham.pptx')
links_path = path.join(
    path.dirname(path.abspath(__file__)),
    '..', 'data', 'video_links.json'
)
graphics_dir = path.join(
    path.dirname(path.abspath(__file__)),
    '..', 'data', 'graphics'
)
media_dir = path.join(
    path.dirname(path.abspath(__file__)),
    '..', 'data', 'nethriq_media'
)

df = pd.read_csv(csv_path)
row = df[df["player_id"] == 0].iloc[0]

with open(links_path, "r") as f:
    video_links = json.load(f)

# ============================================================================
# Load PowerPoint template and create token replacement map
# ============================================================================
prs = Presentation(ppt_template_path)
player_key = f"player_{int(row['player_id'])}"

# Calculate overall grade from average of individual grades
grade_to_num = {"Beginner": 0, "Intermediate": 1, "Advanced": 2, "Pro": 3}
num_to_grade = {0: "Beginner", 1: "Intermediate", 2: "Advanced", 3: "Pro"}

grade_values = [
    grade_to_num.get(row["serve_depth_grade"], 0),
    grade_to_num.get(row["serve_height_grade"], 0),
    grade_to_num.get(row["serve_kitchen_grade"], 0),
    grade_to_num.get(row["return_depth_grade"], 0),
    grade_to_num.get(row["return_height_grade"], 0),
    grade_to_num.get(row["return_kitchen_grade"], 0),
]
avg_grade_num = round(sum(grade_values) / len(grade_values))
overall_grade = num_to_grade.get(avg_grade_num, "Intermediate")

token_map = {
    "{{PLAYER}}": row["player_name"] if pd.notna(row["player_name"]) else "Player",
    "{{SERVE_DEPTH_VALUE}}": f"{row['serve_depth_avg']:.2f}",
    "{{SERVE_DEPTH_GRADE}}": row["serve_depth_grade"],
    "{{SERVE_HEIGHT_VALUE}}": f"{row['serve_height_avg']:.2f}",
    "{{SERVE_HEIGHT_GRADE}}": row["serve_height_grade"],
    "{{RETURN_DEPTH_VALUE}}": f"{row['return_depth_avg']:.2f}",
    "{{RETURN_DEPTH_GRADE}}": row["return_depth_grade"],
    "{{RETURN_HEIGHT_VALUE}}": f"{row['return_height_avg']:.2f}",
    "{{RETURN_HEIGHT_GRADE}}": row["return_height_grade"],
    "{{KAS}}": f"{row['serve_kitchen_pct'] * 100:.0f}",
    "{{KAS_GRADE}}": row["serve_kitchen_grade"],
    "{{KAR}}": f"{row['return_kitchen_pct'] * 100:.0f}",
    "{{KAR_GRADE}}": row["return_kitchen_grade"],
    "{{OVERALL_GRADE}}": overall_grade,
    "{{BEST_SHOTS_VIDEO_LINK}}": "Best Shots",
    "{{RETURN_VIDEO_LINK}}": "Returns in Game",
    "{{SERVE_VIDEO_LINK}}": "Serves in Game"
}

token_map = {key: str(value) for key, value in token_map.items()}

link_map = {
    "{{BEST_SHOTS_VIDEO_LINK}}": video_links.get(
        f"{player_key}_best_shots", {}
    ).get("link"),
    "{{RETURN_VIDEO_LINK}}": video_links.get(
        f"{player_key}_return_context", {}
    ).get("link"),
    "{{SERVE_VIDEO_LINK}}": video_links.get(
        f"{player_key}_serve_context", {}
    ).get("link")
}

# ============================================================================
# Replace tokens in all slides
# ============================================================================
for slide in prs.slides:
    for shape in slide.shapes:
        replace_tokens_and_links(shape, token_map, link_map)

# ============================================================================
# Inject kitchen snapshot image
# ============================================================================
inject_kitchen_snapshot(prs, row['player_id'], graphics_dir)
if THUMBNAIL == "static":
    inject_thumbnail(prs, row['player_id'], media_dir)

# ============================================================================
# Save output PowerPoint file
# ============================================================================
output_dir = path.join(
    path.dirname(path.abspath(__file__)),
    '..',
    'data',
    'delivery_staging',
    f"Player_{int(row['player_id'])}",
    'Reports'
)
os.makedirs(output_dir, exist_ok=True)
output_ppt_path = path.join(output_dir, 'player_report.pptx')
prs.save(output_ppt_path)
print(f"Generated {output_ppt_path}")